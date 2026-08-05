from io import BytesIO

from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject
import pytest

from text_extraction import ExtractionError, extract_text


def test_extract_text_file_cleans_whitespace():
    data = b"  Topic   One\n\n\nDetails\twith   extra spaces  "

    text = extract_text("notes.md", "text/markdown", data)

    assert text == "Topic One\n\nDetails with extra spaces"


def test_extract_docx_reads_paragraphs_and_table_cells():
    document = Document()
    document.add_paragraph("Biology Midterm")
    document.add_paragraph("DNA replication copies genetic material.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Helicase"
    data = _save_docx(document)

    text = extract_text(
        "biology.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        data,
    )

    assert "Biology Midterm" in text
    assert "DNA replication copies genetic material." in text
    assert "Term" in text
    assert "Helicase" in text


def test_extract_pptx_reads_slide_text():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "OSI Model"
    box.text_frame.add_paragraph().text = "Seven layers organize networking."
    data = _save_pptx(presentation)

    text = extract_text(
        "networking.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        data,
    )

    assert text == "OSI Model\n\nSeven layers organize networking."


def test_extract_pdf_removes_repeated_headers_and_page_numbers():
    data = _pdf_bytes(
        "Study Studio\nRouting Basics\nRouters choose paths.\n1",
        "Study Studio\nSubnetting\nSubnets divide networks.\n2",
    )

    text = extract_text("networking.pdf", "application/pdf", data)

    assert "Study Studio" not in text
    assert "Routing Basics" in text
    assert "Routers choose paths." in text
    assert "Subnetting" in text
    assert "Subnets divide networks." in text
    assert "\n\n1" not in text
    assert "\n\n2" not in text


def test_unsupported_file_type_raises_clear_error():
    with pytest.raises(ExtractionError, match="Unsupported file type"):
        extract_text("archive.zip", "application/zip", b"zip bytes")


def test_corrupt_pdf_raises_clear_error():
    with pytest.raises(ExtractionError, match="Could not read broken.pdf"):
        extract_text("broken.pdf", "application/pdf", b"not a real pdf")


def _save_docx(document: Document) -> bytes:
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _save_pptx(presentation: Presentation) -> bytes:
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _pdf_bytes(*pages: str) -> bytes:
    writer = PdfWriter()

    for page_text in pages:
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject(
            {
                NameObject("/Font"): DictionaryObject(
                    {
                        NameObject("/F1"): DictionaryObject(
                            {
                                NameObject("/Type"): NameObject("/Font"),
                                NameObject("/Subtype"): NameObject("/Type1"),
                                NameObject("/BaseFont"): NameObject("/Helvetica"),
                            }
                        )
                    }
                )
            }
        )
        stream = DecodedStreamObject()
        stream.set_data(_pdf_text_stream(page_text))
        page[NameObject("/Contents")] = stream

    buffer = BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def _pdf_text_stream(text: str) -> bytes:
    lines = text.splitlines()
    commands = ["BT", "/F1 12 Tf", "14 TL", "72 720 Td"]
    for line in lines:
        commands.append(f"({_escape_pdf_text(line)}) Tj")
        commands.append("T*")
    commands.append("ET")
    return "\n".join(commands).encode("latin-1")


def _escape_pdf_text(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
