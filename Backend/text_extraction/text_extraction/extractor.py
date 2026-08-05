from __future__ import annotations

from collections import Counter
from io import BytesIO
from pathlib import Path
import re

from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError


class ExtractionError(ValueError):
    pass


_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PAGE_NUMBER_RE = re.compile(r"^(?:page\s*)?\d+(?:\s*(?:of|/)\s*\d+)?$", re.I)
_SPACE_RE = re.compile(r"[ \t\f\v]+")


def extract_text(filename: str, mime: str, data: bytes) -> str:
    kind = _detect_kind(filename, mime)

    try:
        if kind == "text":
            return _clean_blocks([data.decode("utf-8", errors="replace")])
        if kind == "pdf":
            return _extract_pdf(data, filename)
        if kind == "docx":
            return _extract_docx(data)
        if kind == "pptx":
            return _extract_pptx(data)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Could not extract text from {filename}.") from exc

    raise ExtractionError(
        f"Unsupported file type for {filename}. "
        "Supported types: PDF, DOCX, PPTX, TXT, MD."
    )


def _detect_kind(filename: str, mime: str) -> str:
    ext = Path(filename or "").suffix.lower()
    normalized_mime = (mime or "").split(";")[0].strip().lower()

    if normalized_mime.startswith("text/") or ext in {".txt", ".md"}:
        return "text"
    if normalized_mime == "application/pdf" or ext == ".pdf":
        return "pdf"
    if normalized_mime == _DOCX_MIME or ext == ".docx":
        return "docx"
    if normalized_mime == _PPTX_MIME or ext == ".pptx":
        return "pptx"
    return "unsupported"


def _extract_pdf(data: bytes, filename: str) -> str:
    try:
        reader = PdfReader(BytesIO(data))
    except PdfReadError as exc:
        raise ExtractionError(
            f"Could not read {filename}. The file may be corrupt."
        ) from exc

    if reader.is_encrypted:
        try:
            decrypted = reader.decrypt("")
        except Exception as exc:
            raise ExtractionError(
                f"Could not read {filename}. The PDF is encrypted."
            ) from exc
        if decrypted == 0:
            raise ExtractionError(f"Could not read {filename}. The PDF is encrypted.")

    pages = [page.extract_text() or "" for page in reader.pages]
    return _clean_pages(pages)


def _extract_docx(data: bytes) -> str:
    document = Document(BytesIO(data))
    blocks: list[str] = [paragraph.text for paragraph in document.paragraphs]

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                blocks.extend(paragraph.text for paragraph in cell.paragraphs)

    return _clean_blocks(blocks)


def _extract_pptx(data: bytes) -> str:
    presentation = Presentation(BytesIO(data))
    blocks: list[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = " ".join(run.text for run in paragraph.runs)
                blocks.append(text)

    return _clean_blocks(blocks)


def _clean_pages(pages: list[str]) -> str:
    page_lines = [_clean_lines(page) for page in pages]
    repeated = _repeated_edge_lines(page_lines)
    blocks: list[str] = []

    for lines in page_lines:
        for line in lines:
            key = _line_key(line)
            if key in repeated or _PAGE_NUMBER_RE.match(line):
                continue
            blocks.append(line)

    return "\n\n".join(blocks)


def _clean_blocks(blocks: list[str]) -> str:
    lines: list[str] = []
    for block in blocks:
        lines.extend(_clean_lines(block))
    return "\n\n".join(lines)


def _clean_lines(text: str) -> list[str]:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n").replace("\xa0", " ")
    lines: list[str] = []

    for line in normalized.split("\n"):
        cleaned = _SPACE_RE.sub(" ", line).strip()
        if cleaned:
            lines.append(cleaned)

    return lines


def _repeated_edge_lines(page_lines: list[list[str]]) -> set[str]:
    if len(page_lines) < 2:
        return set()

    counts: Counter[str] = Counter()
    for lines in page_lines:
        edge_lines = {_line_key(line) for line in [*lines[:2], *lines[-2:]] if line}
        counts.update(edge_lines)

    threshold = max(2, (len(page_lines) + 1) // 2)
    return {
        line
        for line, count in counts.items()
        if count >= threshold and len(line) <= 100
    }


def _line_key(line: str) -> str:
    return _SPACE_RE.sub(" ", line).strip().lower()
